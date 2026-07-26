import logging
from pathlib import Path

from app.celery_app import celery_app
from app.config import get_settings

logger = logging.getLogger(__name__)


@celery_app.task(bind=True, name="app.tasks.pdf_to_word")
def pdf_to_word(self, input_path: str, original_filename: str = "document.pdf", mode: str = "layout") -> dict:
    settings = get_settings()
    stem = Path(original_filename).stem
    output_path = settings.RESULT_DIR / f"{self.request.id}.docx"
    
    logger.info("Converting %s to DOCX in %s mode", original_filename, mode)
    
    try:
        if mode == "flowing":
            import docx
            import pymupdf
            from docx.shared import Inches, Pt, RGBColor
            
            pdf_in = pymupdf.open(input_path)
            doc_out = docx.Document()
            
            # Apply exact geometry from first page
            if len(pdf_in) > 0:
                first_page = pdf_in[0]
                rect = first_page.rect
                page_width_in = rect.width / 72.0
                page_height_in = rect.height / 72.0
                
                blocks = first_page.get_text("blocks")
                valid_blocks = [b for b in blocks if len(b) > 4 and str(b[4]).strip()]
                if valid_blocks:
                    min_x0 = min(b[0] for b in valid_blocks)
                    min_y0 = min(b[1] for b in valid_blocks)
                    max_x1 = max(b[2] for b in valid_blocks)
                    max_y1 = max(b[3] for b in valid_blocks)
                    top_m = max(0.2, min_y0 / 72.0 - 0.05)
                    bot_m = max(0.2, (rect.height - max_y1) / 72.0 - 0.05)
                    left_m = max(0.2, min_x0 / 72.0 - 0.05)
                    right_m = max(0.2, (rect.width - max_x1) / 72.0 - 0.05)
                else:
                    top_m = bot_m = left_m = right_m = 0.35

                for section in doc_out.sections:
                    section.page_width = Inches(page_width_in)
                    section.page_height = Inches(page_height_in)
                    section.top_margin = Inches(top_m)
                    section.bottom_margin = Inches(bot_m)
                    section.left_margin = Inches(left_m)
                    section.right_margin = Inches(right_m)

            total = len(pdf_in)
            for i, page in enumerate(pdf_in):
                blocks = page.get_text("dict").get("blocks", [])
                for b in blocks:
                    if b.get("type") == 0:  # text block
                        p = doc_out.add_paragraph()
                        p.paragraph_format.space_before = Pt(0)
                        p.paragraph_format.space_after = Pt(1)
                        p.paragraph_format.line_spacing = 1.0
                        for line in b.get("lines", []):
                            for span in line.get("spans", []):
                                text = span.get("text", "")
                                if text:
                                    run = p.add_run(text)
                                    run.font.size = Pt(span.get("size", 11))
                                    flags = span.get("flags", 0)
                                    if flags & 16:
                                        run.font.bold = True
                                    if flags & 2:
                                        run.font.italic = True
                                    color = span.get("color", 0)
                                    if color != 0:
                                        try:
                                            run.font.color.rgb = RGBColor.from_string(f"{color:06x}")
                                        except Exception:
                                            pass
                            p.add_run(" ")
                
                progress = int(((i + 1) / total) * 100)
                self.update_state(state="PROGRESS", meta={"progress": progress, "filename": f"{stem}.docx"})
                
            doc_out.save(str(output_path))
            pdf_in.close()
            
        else:
            import docx
            import pymupdf
            from docx.shared import Inches, Pt
            from pdf2docx import Converter
            
            cv = Converter(input_path)
            self.update_state(state="PROGRESS", meta={"progress": 25, "filename": f"{stem}.docx"})
            cv.convert(str(output_path), start=0, end=None)
            cv.close()
            
            # Post-process generated DOCX to match exact PDF page margins & prevent page overflow
            try:
                pdf_in = pymupdf.open(input_path)
                if len(pdf_in) > 0:
                    first_page = pdf_in[0]
                    rect = first_page.rect
                    page_width_in = rect.width / 72.0
                    page_height_in = rect.height / 72.0
                    
                    blocks = first_page.get_text("blocks")
                    valid_blocks = [b for b in blocks if len(b) > 4 and str(b[4]).strip()]
                    if valid_blocks:
                        min_x0 = min(b[0] for b in valid_blocks)
                        min_y0 = min(b[1] for b in valid_blocks)
                        max_x1 = max(b[2] for b in valid_blocks)
                        max_y1 = max(b[3] for b in valid_blocks)
                        top_m = max(0.2, min_y0 / 72.0 - 0.05)
                        bot_m = max(0.2, (rect.height - max_y1) / 72.0 - 0.05)
                        left_m = max(0.2, min_x0 / 72.0 - 0.05)
                        right_m = max(0.2, (rect.width - max_x1) / 72.0 - 0.05)
                    else:
                        top_m = bot_m = left_m = right_m = 0.35

                    doc_mod = docx.Document(str(output_path))
                    for section in doc_mod.sections:
                        section.page_width = Inches(page_width_in)
                        section.page_height = Inches(page_height_in)
                        section.top_margin = Inches(top_m)
                        section.bottom_margin = Inches(bot_m)
                        section.left_margin = Inches(left_m)
                        section.right_margin = Inches(right_m)

                    for p in doc_mod.paragraphs:
                        p.paragraph_format.space_before = Pt(0)
                        p.paragraph_format.space_after = Pt(1)
                        p.paragraph_format.line_spacing = 1.0

                    for t in doc_mod.tables:
                        for row in t.rows:
                            for cell in row.cells:
                                for p in cell.paragraphs:
                                    p.paragraph_format.space_before = Pt(0)
                                    p.paragraph_format.space_after = Pt(1)
                                    p.paragraph_format.line_spacing = 1.0

                    doc_mod.save(str(output_path))
                pdf_in.close()
            except Exception as e:
                logger.warning(f"Could not post-process DOCX geometry: {e}")

            self.update_state(state="PROGRESS", meta={"progress": 100, "filename": f"{stem}.docx"})
            
        return {"result_path": str(output_path), "filename": f"{stem}.docx"}
    except Exception:
        logger.exception("Error converting PDF to Word")
        raise


@celery_app.task(bind=True, name="app.tasks.pdf_to_powerpoint")
def pdf_to_powerpoint(self, input_path: str, original_filename: str = "document.pdf") -> dict:
    settings = get_settings()
    stem = Path(original_filename).stem
    output_path = settings.RESULT_DIR / f"{self.request.id}.pptx"
    
    logger.info("Converting %s to PPTX", original_filename)
    
    try:
        import pymupdf
        from pptx import Presentation
        from pptx.util import Inches
        
        doc = pymupdf.open(input_path)
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6] 
        
        total = len(doc)
        if total == 0:
            raise ValueError("PDF has no pages")
            
        for i in range(total):
            page = doc[i]
            pix = page.get_pixmap(dpi=150)
            temp_img_path = settings.RESULT_DIR / f"{self.request.id}_page_{i}.png"
            pix.save(str(temp_img_path))
            
            if i == 0:
                prs.slide_width = int(page.rect.width * 12700)
                prs.slide_height = int(page.rect.height * 12700)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(str(temp_img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            temp_img_path.unlink()
            
            progress = int(((i + 1) / total) * 100)
            self.update_state(state="PROGRESS", meta={"progress": progress, "filename": f"{stem}.pptx"})
            
        prs.save(str(output_path))
        doc.close()
        
        return {"result_path": str(output_path), "filename": f"{stem}.pptx"}
    except Exception:
        logger.exception("Error converting PDF to PPTX")
        raise


@celery_app.task(bind=True, name="app.tasks.pdf_to_excel")
def pdf_to_excel(self, input_path: str, original_filename: str = "document.pdf") -> dict:
    settings = get_settings()
    stem = Path(original_filename).stem
    output_path = settings.RESULT_DIR / f"{self.request.id}.xlsx"
    
    logger.info("Converting %s to XLSX", original_filename)
    
    try:
        import pdfplumber
        from openpyxl import Workbook
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Extracted Tables"
        
        row_offset = 1
        with pdfplumber.open(input_path) as pdf:
            total = len(pdf.pages)
            if total == 0:
                raise ValueError("PDF has no pages")
                
            for i, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        clean_row = [str(cell).strip() if cell is not None else "" for cell in row]
                        for col_idx, cell_value in enumerate(clean_row, start=1):
                            ws.cell(row=row_offset, column=col_idx, value=cell_value)
                        row_offset += 1
                    # Add a blank row between tables
                    row_offset += 1
                
                progress = int(((i + 1) / total) * 100)
                self.update_state(state="PROGRESS", meta={"progress": progress, "filename": f"{stem}.xlsx"})
                
        # If no tables found, at least return an empty excel sheet rather than failing completely
        wb.save(str(output_path))
        
        return {"result_path": str(output_path), "filename": f"{stem}.xlsx"}
    except Exception:
        logger.exception("Error converting PDF to Excel")
        raise
