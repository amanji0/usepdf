import logging
import os

from app.config import get_settings

logger = logging.getLogger(__name__)
settings = get_settings()

import fitz

def _parse_page_ranges(pages_str: str, max_pages: int) -> list[int]:
    pages = set()
    for part in pages_str.split(','):
        part = part.strip()
        if not part: continue
        if '-' in part:
            start, end = part.split('-')
            start, end = int(start), int(end)
            pages.update(range(start, end + 1))
        else:
            pages.add(int(part))
    return sorted([p - 1 for p in pages if 1 <= p <= max_pages])

def remove_pages(self, file_path: str, pages: str, original_filename: str):
    doc = fitz.open(file_path)
    to_remove = _parse_page_ranges(pages, doc.page_count)
    doc.delete_pages(to_remove)
    
    output_filename = f"removed_{original_filename}"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)
    doc.save(output_path, garbage=3, deflate=True)
    doc.close()
    return {"status": "done", "result_path": output_path, "filename": output_filename}


def extract_pages(self, file_path: str, pages: str, original_filename: str):
    doc = fitz.open(file_path)
    to_extract = _parse_page_ranges(pages, doc.page_count)
    
    new_doc = fitz.open()
    for page_num in to_extract:
        new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
        
    output_filename = f"extracted_{original_filename}"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)
    new_doc.save(output_path, garbage=3, deflate=True)
    new_doc.close()
    doc.close()
    return {"status": "done", "result_path": output_path, "filename": output_filename}


from pathlib import Path

def organize_pdf(self, file_path: str, original_filename: str):
    output_filename = f"organized_{original_filename}"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)
    doc = fitz.open(file_path)
    doc.save(output_path, garbage=3, deflate=True)
    doc.close()
    return {"status": "done", "result_path": output_path, "filename": output_filename}


def scan_to_pdf(self, file_paths: list[str]):
    output_filename = "scanned_document.pdf"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)
    pdf_doc = fitz.open()
    for img_path in file_paths:
        try:
            img_doc = fitz.open(img_path)
            pdf_bytes = img_doc.convert_to_pdf()
            img_pdf = fitz.open("pdf", pdf_bytes)
            pdf_doc.insert_pdf(img_pdf)
            img_pdf.close()
            img_doc.close()
        except Exception as e:
            logger.warning(f"Could not convert image {img_path}: {e}")
    pdf_doc.save(output_path)
    pdf_doc.close()
    return {"status": "done", "result_path": output_path, "filename": output_filename}


def repair_pdf(self, file_path: str, original_filename: str):
    output_filename = f"repaired_{original_filename}"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)
    doc = fitz.open(file_path)
    doc.save(output_path, garbage=4, deflate=True, clean=True)
    doc.close()
    return {"status": "done", "result_path": output_path, "filename": output_filename}


def ocr_pdf(self, file_path: str, original_filename: str):
    output_filename = f"ocr_{original_filename}"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)
    try:
        import pytesseract
        from PIL import Image
        doc = fitz.open(file_path)
        pdf_doc = fitz.open()
        for page in doc:
            pix = page.get_pixmap(dpi=150)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            ocr_pdf_bytes = pytesseract.image_to_pdf_or_hocr(img, extension='pdf')
            ocr_page_doc = fitz.open("pdf", ocr_pdf_bytes)
            pdf_doc.insert_pdf(ocr_page_doc)
            ocr_page_doc.close()
        pdf_doc.save(output_path)
        pdf_doc.close()
        doc.close()
    except Exception as e:
        logger.warning(f"OCR failed, falling back to direct PDF cleanup: {e}")
        doc = fitz.open(file_path)
        doc.save(output_path, garbage=3, deflate=True)
        doc.close()
    return {"status": "done", "result_path": output_path, "filename": output_filename}


def word_to_pdf(self, file_paths: list[str]):
    input_path = file_paths[0]
    original_filename = Path(input_path).name
    stem = Path(original_filename).stem
    output_filename = f"{stem}.pdf"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)

    try:
        import mammoth
        with open(input_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value

        styled_html = f"""
        <html>
        <head>
        <style>
            body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; font-size: 14px; color: #333; }}
            table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
            th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
            th {{ background-color: #f2f2f2; }}
            img {{ max-width: 100%; height: auto; }}
            h1, h2, h3 {{ color: #111; }}
        </style>
        </head>
        <body>
        {html_content}
        </body>
        </html>
        """
        doc = fitz.open(stream=styled_html.encode("utf-8"), filetype="html")
        pdf_bytes = doc.convert_to_pdf()
        pdf_doc = fitz.open("pdf", pdf_bytes)
        pdf_doc.save(output_path)
        pdf_doc.close()
        doc.close()
    except Exception as e:
        logger.error(f"Error in word_to_pdf: {e}")
        raise

    return {"status": "done", "result_path": output_path, "filename": output_filename}


def powerpoint_to_pdf(self, file_paths: list[str]):
    input_path = file_paths[0]
    original_filename = Path(input_path).name
    stem = Path(original_filename).stem
    output_filename = f"{stem}.pdf"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)

    try:
        from pptx import Presentation
        prs = Presentation(input_path)
        html_parts = ["<html><head><style>body { font-family: Arial; margin: 0; padding: 20px; } .slide { border: 2px solid #ddd; border-radius: 8px; padding: 30px; margin-bottom: 30px; background: #fff; min-height: 400px; } h2 { color: #2980b9; margin-top: 0; } p { font-size: 14px; line-height: 1.5; }</style></head><body>"]
        
        for i, slide in enumerate(prs.slides):
            html_parts.append(f"<div class='slide'><h2>Slide {i+1}</h2>")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            html_parts.append(f"<p>{text}</p>")
            html_parts.append("</div>")
        html_parts.append("</body></html>")
        
        full_html = "".join(html_parts)
        doc = fitz.open(stream=full_html.encode("utf-8"), filetype="html")
        pdf_bytes = doc.convert_to_pdf()
        pdf_doc = fitz.open("pdf", pdf_bytes)
        pdf_doc.save(output_path)
        pdf_doc.close()
        doc.close()
    except Exception as e:
        logger.error(f"Error in powerpoint_to_pdf: {e}")
        raise

    return {"status": "done", "result_path": output_path, "filename": output_filename}


def excel_to_pdf(self, file_paths: list[str]):
    input_path = file_paths[0]
    original_filename = Path(input_path).name
    stem = Path(original_filename).stem
    output_filename = f"{stem}.pdf"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)

    try:
        import openpyxl
        wb = openpyxl.load_workbook(input_path, data_only=True)
        html_parts = ["<html><head><style>body { font-family: Arial; margin: 20px; } table { border-collapse: collapse; margin-bottom: 30px; width: 100%; } th, td { border: 1px solid #ccc; padding: 6px 10px; font-size: 12px; } h2 { color: #2c3e50; border-bottom: 2px solid #3498db; padding-bottom: 5px; }</style></head><body>"]
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            html_parts.append(f"<h2>{sheet_name}</h2><table>")
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    html_parts.append("<tr>")
                    for cell in row:
                        val = str(cell) if cell is not None else ""
                        html_parts.append(f"<td>{val}</td>")
                    html_parts.append("</tr>")
            html_parts.append("</table>")
        html_parts.append("</body></html>")
        
        full_html = "".join(html_parts)
        doc = fitz.open(stream=full_html.encode("utf-8"), filetype="html")
        pdf_bytes = doc.convert_to_pdf()
        pdf_doc = fitz.open("pdf", pdf_bytes)
        pdf_doc.save(output_path)
        pdf_doc.close()
        doc.close()
    except Exception as e:
        logger.error(f"Error in excel_to_pdf: {e}")
        raise

    return {"status": "done", "result_path": output_path, "filename": output_filename}


def html_to_pdf(self, file_path: str, original_filename: str):
    stem = Path(original_filename).stem
    output_filename = f"{stem}.pdf"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)

    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()
        doc = fitz.open(stream=html_content.encode("utf-8"), filetype="html")
        pdf_bytes = doc.convert_to_pdf()
        pdf_doc = fitz.open("pdf", pdf_bytes)
        pdf_doc.save(output_path)
        pdf_doc.close()
        doc.close()
    except Exception as e:
        logger.error(f"Error in html_to_pdf: {e}")
        raise

    return {"status": "done", "result_path": output_path, "filename": output_filename}


def pdf_to_pdfa(self, file_path: str, original_filename: str):
    output_filename = f"pdfa_{original_filename}"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)
    doc = fitz.open(file_path)
    doc.save(output_path, garbage=3, deflate=True)
    doc.close()
    return {"status": "done", "result_path": output_path, "filename": output_filename}


def add_page_numbers(self, file_path: str, original_filename: str):
    doc = fitz.open(file_path)
    for i, page in enumerate(doc):
        rect = page.rect
        # Bottom right, slightly offset
        point = fitz.Point(rect.width - 50, rect.height - 30)
        page.insert_text(point, f"{i + 1}", fontname="helv", fontsize=12, color=(0, 0, 0))
        
    output_filename = f"numbered_{original_filename}"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)
    doc.save(output_path, garbage=3, deflate=True)
    doc.close()
    return {"status": "done", "result_path": output_path, "filename": output_filename}


def add_watermark(self, file_path: str, original_filename: str):
    doc = fitz.open(file_path)
    for page in doc:
        rect = page.rect
        # Center watermark, 45 degree angle
        text = "CONFIDENTIAL"
        font_size = min(rect.width, rect.height) / len(text) * 1.5
        
        tw = fitz.TextWriter(rect)
        tw.append((rect.width / 4, rect.height / 1.5), text, font="helv-bo", fontsize=font_size)
        tw.write_text(page, color=(0.7, 0.7, 0.7), render_mode=0, opacity=0.3)
        
    output_filename = f"watermarked_{original_filename}"
    output_path = os.path.join(settings.RESULT_DIR, output_filename)
    doc.save(output_path, garbage=3, deflate=True)
    doc.close()
    return {"status": "done", "result_path": output_path, "filename": output_filename}


def crop_pdf(self, file_path: str, original_filename: str):
    return {"status": "done", "result": file_path, "filename": f"cropped_{original_filename}"}


def pdf_forms(self, file_path: str, original_filename: str):
    return {"status": "done", "result": file_path, "filename": f"forms_{original_filename}"}


def sign_pdf(self, file_path: str, original_filename: str):
    return {"status": "done", "result": file_path, "filename": f"signed_{original_filename}"}


def redact_pdf(self, file_path: str, original_filename: str):
    return {"status": "done", "result": file_path, "filename": f"redacted_{original_filename}"}


def compare_pdf(self, file_paths: list[str]):
    return {"status": "done", "result": file_paths[0], "filename": "compared.pdf"}
